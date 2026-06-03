"""PowerPoint .pptx reader via python-pptx."""

from pathlib import Path

MAX_CHARS = 50_000


def read(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        # Slide notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"--- Slide {i} notes ---\n{notes}")
    out = "\n\n".join(parts)
    if len(out) > MAX_CHARS:
        out = out[:MAX_CHARS] + f"\n\n... (truncated at {MAX_CHARS} chars)"
    return out or "(empty presentation)"
